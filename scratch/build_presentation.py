import sys
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Colors
    C_BG_DARK = RGBColor(9, 19, 25)         # #091319
    C_BG_DARK_CARD = RGBColor(15, 36, 45)   # #0f242d
    C_TEAL_ACCENT = RGBColor(20, 184, 166)  # #14b8a6
    C_TEAL_DARK = RGBColor(13, 148, 136)    # #0d9488
    C_TEAL_LIGHT = RGBColor(204, 251, 241)  # #ccfbf1
    C_BG_LIGHT = RGBColor(248, 250, 252)    # #f8fafc
    C_CARD_LIGHT = RGBColor(255, 255, 255)  # #ffffff
    C_CARD_BORDER = RGBColor(226, 232, 240) # #e2e8f0
    C_TEXT_DARK = RGBColor(15, 23, 42)      # #0f172a
    C_TEXT_MUTED = RGBColor(71, 85, 105)    # #475569
    C_TEXT_LIGHT = RGBColor(255, 255, 255)
    C_TEXT_LIGHT_MUTED = RGBColor(148, 163, 184)
    C_AMBER = RGBColor(245, 158, 11)        # Warning / highlight
    C_BLUE_BADGE = RGBColor(30, 58, 138)
    
    blank_slide_layout = prs.slide_layouts[6]
    
    logo_path = "/Users/dendyaditya/Projects/windy_project/public/images/nsg-logo.jpg"
    has_logo = os.path.exists(logo_path)
    
    def add_header(slide, title_text, category_text="PT NATTU GLOBAL SYNERGY", dark=False):
        # Category / Overline
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(10), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = C_TEAL_ACCENT if dark else C_TEAL_DARK
        
        # Title
        t_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.75), Inches(9.5), Inches(0.6))
        tf_t = t_box.text_frame
        tf_t.word_wrap = True
        tf_t.margin_left = tf_t.margin_top = tf_t.margin_right = tf_t.margin_bottom = 0
        p_t = tf_t.paragraphs[0]
        p_t.text = title_text
        p_t.font.size = Pt(22)
        p_t.font.bold = True
        p_t.font.color.rgb = C_TEXT_LIGHT if dark else C_TEXT_DARK
        
        # Top right logo if exists
        if has_logo:
            try:
                slide.shapes.add_picture(logo_path, Inches(11.2), Inches(0.4), height=Inches(0.8))
            except:
                pass

    def add_card(slide, x, y, w, h, bg_color=C_CARD_LIGHT, border_color=C_CARD_BORDER):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        if border_color:
            shape.line.color.rgb = border_color
            shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        return shape

    # -------------------------------------------------------------
    # SLIDE 1: COVER (Dark Theme)
    # -------------------------------------------------------------
    s1 = prs.slides.add_slide(blank_slide_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = C_BG_DARK
    bg1.line.fill.background()
    
    # Decorative accent card
    dec = s1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(11.73), Inches(5.5))
    dec.fill.solid()
    dec.fill.fore_color.rgb = C_BG_DARK_CARD
    dec.line.color.rgb = C_TEAL_DARK
    dec.line.width = Pt(1.5)
    
    if has_logo:
        try:
            s1.shapes.add_picture(logo_path, Inches(1.3), Inches(1.5), height=Inches(1.2))
        except:
            pass
            
    tb_cover = s1.shapes.add_textbox(Inches(1.3), Inches(2.9), Inches(10.5), Inches(3.2))
    tf_c = tb_cover.text_frame
    tf_c.word_wrap = True
    
    p1 = tf_c.paragraphs[0]
    p1.text = "PROJECT HANDOVER & USER TRAINING"
    p1.font.size = Pt(12)
    p1.font.bold = True
    p1.font.color.rgb = C_TEAL_ACCENT
    p1.space_after = Pt(10)
    
    p2 = tf_c.add_paragraph()
    p2.text = "Website Company Profile &\nIntegrasi Custom Email Gmail"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = C_TEXT_LIGHT
    p2.space_after = Pt(16)
    
    p3 = tf_c.add_paragraph()
    p3.text = "PT Nattu Global Synergy  |  Domain: nattuglobalsynergy.co.id\nPanduan Penggunaan, Integrasi Third-Party Mail Client, & SOP Layanan Support"
    p3.font.size = Pt(14)
    p3.font.color.rgb = C_TEXT_LIGHT_MUTED

    # -------------------------------------------------------------
    # SLIDE 2: AGENDA & EXECUTIVE SUMMARY (Light Theme)
    # -------------------------------------------------------------
    s2 = prs.slides.add_slide(blank_slide_layout)
    add_header(s2, "Agenda & Ringkasan Hasil Proyek")
    
    # 3 Summary Cards
    cards_data_s2 = [
        ("01. Website Company Profile", [
            "• Website resmi bilingual (ID & EN) live & responsive",
            "• Desain modern, cepat (Next.js static), SEO ready",
            "• Formulir kontak terintegrasi langsung ke email resmi",
            "• Showcase 3 divisi bisnis utama & portofolio sektor"
        ], C_TEAL_LIGHT, C_TEAL_DARK),
        ("02. Integrasi Domain Email di Gmail", [
            "• Akun email resmi domain @nattuglobalsynergy.co.id",
            "• Integrasi POP3 (terima) & SMTP (kirim) via Gmail",
            "• Bisa diakses via Web Browser & Gmail Mobile App",
            "• Tanpa lisensi Google Workspace berbayar bulanan"
        ], RGBColor(238, 242, 255), RGBColor(79, 70, 229)),
        ("03. SOP Layanan & Dukungan", [
            "• Prosedur pengajuan penambahan email baru",
            "• Tata cara reset password & penanganan kendala",
            "• Alur permintaan update konten & bug fixing website",
            "• SLA (Service Level Agreement) pengerjaan & support"
        ], RGBColor(254, 243, 199), RGBColor(217, 119, 6)),
    ]
    
    for i, (title, items, bg_c, acc_c) in enumerate(cards_data_s2):
        x = 0.8 + (i * 3.95)
        add_card(s2, x, 1.6, 3.8, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
        
        # Header strip inside card
        header_shape = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(1.75), Inches(3.5), Inches(0.65))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = bg_c
        header_shape.line.fill.background()
        tf = header_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = acc_c
        p.alignment = PP_ALIGN.CENTER
        
        # Content
        tb = s2.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(3.4), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(12)
            p.font.color.rgb = C_TEXT_MUTED
            p.space_after = Pt(10)

    # -------------------------------------------------------------
    # SLIDE 3: OVERVIEW WEBSITE & FITUR UTAMA
    # -------------------------------------------------------------
    s3 = prs.slides.add_slide(blank_slide_layout)
    add_header(s3, "Struktur & Fitur Website PT Nattu Global Synergy", "MODUL 1: WEBSITE DEVELOPMENT")
    
    add_card(s3, 0.8, 1.6, 5.6, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_w1 = s3.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_w1 = tb_w1.text_frame
    tf_w1.word_wrap = True
    
    p = tf_w1.paragraphs[0]
    p.text = "Spesifikasi & Arsitektur Website"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(14)
    
    specs = [
        ("URL Domain Resmi", "https://www.nattuglobalsynergy.co.id"),
        ("Teknologi Frontend", "Next.js Static Export + Tailwind CSS"),
        ("Bahasa (Bilingual)", "Bahasa Indonesia (ID) & English (EN)"),
        ("Infrastruktur Hosting", "Biznet Gio Cloud cPanel (SSL Secured)"),
        ("Formulir Kontak", "Web3Forms API (Terkirim langsung ke email)"),
        ("Responsive Design", "Optimal di Smartphone, Tablet, Laptop, & 4K")
    ]
    for label, val in specs:
        p = tf_w1.add_paragraph()
        p.text = f"• {label}: "
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEXT_DARK
        run = p.add_run()
        run.text = val
        run.font.bold = False
        run.font.color.rgb = C_TEXT_MUTED
        p.space_after = Pt(8)

    add_card(s3, 6.8, 1.6, 5.7, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_w2 = s3.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf_w2 = tb_w2.text_frame
    tf_w2.word_wrap = True
    
    p = tf_w2.paragraphs[0]
    p.text = "Seksi Konten & Showcase Bisnis"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(14)
    
    sections = [
        ("Hero & Tagline", "Membangun Fondasi Kuat, Menghadirkan Sinergi Berkelanjutan."),
        ("Tentang Kami & Nilai", "Visi, 5 Misi utama, dan 5 Nilai Perusahaan (Integritas, Kualitas, dll)."),
        ("3 Divisi Bisnis Utama", "1. Konstruksi & Manufaktur\n2. Logam & Material Industri\n3. Komponen Elektronik & Energi Terbarukan"),
        ("Keunggulan & K3", "8 Keunggulan kompetitif & komitmen standar keselamatan kerja K3."),
        ("13 Sektor Industri", "Target industri: Otomotif, Kelistrikan, Petrokimia, Infrastruktur, dll.")
    ]
    for label, desc in sections:
        p = tf_w2.add_paragraph()
        p.text = f"✔ {label}"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEAL_DARK
        p.space_after = Pt(2)
        
        p_sub = tf_w2.add_paragraph()
        p_sub.text = desc
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_TEXT_MUTED
        p_sub.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 4: ARSITEKTUR & KEUNGGULAN INTEGRASI GMAIL
    # -------------------------------------------------------------
    s4 = prs.slides.add_slide(blank_slide_layout)
    add_header(s4, "Arsitektur Integrasi Email Domain via Gmail", "MODUL 2: EMAIL INTEGRATION")
    
    # 3 Columns explaining why and how
    arch_data = [
        ("1. Keamanan & Mail Server", [
            "• Menggunakan cPanel Mail Server Biznet Gio",
            "• Protokol terenkripsi penuh: SSL/TLS",
            "• Kuota penyimpanan terkelola aman",
            "• Domain resmi membangun kredibilitas: @nattuglobalsynergy.co.id"
        ], "Infrastruktur cPanel"),
        ("2. Gmail Client Integration", [
            "• Menggunakan antarmuka Gmail yang familiar",
            "• Bebas biaya lisensi bulanan Google Workspace",
            "• Filter anti-spam Google yang sangat canggih",
            "• Fitur pencarian email cepat & penyimpanan Google Drive"
        ], "Antarmuka Pengguna"),
        ("3. Parameter Teknis Koneksi", [
            "• POP3 Server: mail.nattuglobalsynergy.co.id",
            "• POP3 Port: 995 (SSL/TLS Aktif)",
            "• SMTP Server: mail.nattuglobalsynergy.co.id",
            "• SMTP Port: 465 (SSL/TLS Aktif)",
            "• Username: Alamat email lengkap"
        ], "Konfigurasi Port")
    ]
    
    for i, (title, items, tag) in enumerate(arch_data):
        x = 0.8 + (i * 3.95)
        add_card(s4, x, 1.6, 3.8, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
        
        # Tag
        tb_tag = s4.shapes.add_textbox(Inches(x + 0.2), Inches(1.8), Inches(3.4), Inches(0.3))
        p = tb_tag.text_frame.paragraphs[0]
        p.text = tag.upper()
        p.font.size = Pt(9.5)
        p.font.bold = True
        p.font.color.rgb = C_TEAL_DARK
        
        # Title
        tb_t = s4.shapes.add_textbox(Inches(x + 0.2), Inches(2.1), Inches(3.4), Inches(0.6))
        p = tb_t.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK
        
        # Items
        tb_i = s4.shapes.add_textbox(Inches(x + 0.2), Inches(2.8), Inches(3.4), Inches(3.8))
        tf = tb_i.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11.5)
            p.font.color.rgb = C_TEXT_MUTED
            p.space_after = Pt(8)

    # -------------------------------------------------------------
    # SLIDE 5: STEP-BY-STEP TERIMA EMAIL DI GMAIL (POP3)
    # -------------------------------------------------------------
    s5 = prs.slides.add_slide(blank_slide_layout)
    add_header(s5, "Langkah 1: Menghubungkan Inbox (Menerima Email)", "TUTORIAL SETUP GMAIL")
    
    steps_pop3 = [
        ("Langkah 01", "Buka Pengaturan Gmail", "Login ke akun Gmail (personal/bisnis Anda). Klik ikon roda gigi (Settings) di pojok kanan atas, lalu klik 'See all settings' / 'Lihat semua setelan'."),
        ("Langkah 02", "Pilih Tab Accounts & Import", "Pilih tab 'Accounts and Import'. Cari seksi 'Check mail from other accounts' (Periksa email dari akun lain), lalu klik 'Add an email account'."),
        ("Langkah 03", "Masukkan Email Domain", "Masukkan alamat email domain Anda (contoh: info@nattuglobalsynergy.co.id). Pilih opsi 'Import emails from my other account (POP3)' lalu klik Next."),
        ("Langkah 04", "Input Kredensial Server", "• Username: email lengkap\n• Password: password email\n• POP Server: mail.nattuglobalsynergy.co.id\n• Port: 995\n• Centang: 'Always use a secure connection (SSL)'\n• Centang: 'Label incoming messages'")
    ]
    
    for i, (num, title, desc) in enumerate(steps_pop3):
        x = 0.8 + (i * 2.95)
        add_card(s5, x, 1.6, 2.8, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
        
        # Step badge
        badge = s5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(1.8), Inches(1.2), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_TEAL_DARK
        badge.line.fill.background()
        p = badge.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
        # Title
        tb_t = s5.shapes.add_textbox(Inches(x + 0.2), Inches(2.35), Inches(2.4), Inches(0.6))
        p = tb_t.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK
        
        # Desc
        tb_d = s5.shapes.add_textbox(Inches(x + 0.2), Inches(3.0), Inches(2.4), Inches(3.6))
        tf = tb_d.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 6: STEP-BY-STEP KIRIM EMAIL DARI GMAIL (SMTP)
    # -------------------------------------------------------------
    s6 = prs.slides.add_slide(blank_slide_layout)
    add_header(s6, "Langkah 2: Menghubungkan Outbox (Mengirim Email)", "TUTORIAL SETUP GMAIL")
    
    steps_smtp = [
        ("Langkah 01", "Pilih 'Send mail as'", "Pada tab 'Accounts and Import', cari seksi 'Send mail as' (Kirim email sebagai), lalu klik 'Add another email address' (Tambahkan alamat email lain)."),
        ("Langkah 02", "Isi Nama Pengirim", "Masukkan Nama Perusahaan / Nama Anda (misal: 'PT Nattu Global Synergy - Info') dan alamat email resmi. Centang opsi 'Treat as an alias' lalu klik Next Step."),
        ("Langkah 03", "Konfigurasi Server SMTP", "• SMTP Server: mail.nattuglobalsynergy.co.id\n• Port: 465\n• Username: email lengkap\n• Password: password email\n• Pilih: Secured connection using SSL"),
        ("Langkah 04", "Verifikasi Kode", "Buka inbox Gmail Anda, cari email konfirmasi dari Gmail Team berisi kode verifikasi. Masukkan kode tersebut dan klik 'Verify'. Selesai!")
    ]
    
    for i, (num, title, desc) in enumerate(steps_smtp):
        x = 0.8 + (i * 2.95)
        add_card(s6, x, 1.6, 2.8, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
        
        badge = s6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.2), Inches(1.8), Inches(1.2), Inches(0.4))
        badge.fill.solid()
        badge.fill.fore_color.rgb = RGBColor(79, 70, 229)
        badge.line.fill.background()
        p = badge.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(10)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_LIGHT
        p.alignment = PP_ALIGN.CENTER
        
        tb_t = s6.shapes.add_textbox(Inches(x + 0.2), Inches(2.35), Inches(2.4), Inches(0.6))
        p = tb_t.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_TEXT_DARK
        
        tb_d = s6.shapes.add_textbox(Inches(x + 0.2), Inches(3.0), Inches(2.4), Inches(3.6))
        tf = tb_d.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = C_TEXT_MUTED

    # -------------------------------------------------------------
    # SLIDE 7: CARA MENGGUNAKAN & AKSES MOBILE (SMARTPHONE)
    # -------------------------------------------------------------
    s7 = prs.slides.add_slide(blank_slide_layout)
    add_header(s7, "Panduan Pemakaian Sehari-hari & Akses di Smartphone", "USER GUIDE")
    
    # Left Card: Cara Mengirim & Memilih Sender di Gmail
    add_card(s7, 0.8, 1.6, 5.6, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_u1 = s7.shapes.add_textbox(Inches(1.1), Inches(1.8), Inches(5.0), Inches(4.8))
    tf_u1 = tb_u1.text_frame
    tf_u1.word_wrap = True
    
    p = tf_u1.paragraphs[0]
    p.text = "Cara Mengirim Email Resmi dari Gmail"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(12)
    
    usage_tips = [
        ("1. Klik 'Tulis' (Compose)", "Buka jendela pesan baru di Gmail seperti biasa."),
        ("2. Pilih Dropdown 'From' (Dari)", "Klik baris pengirim, ubah dari email pribadi ke: nama@nattuglobalsynergy.co.id."),
        ("3. Setting Pengirim Default (Opsional)", "Di tab Accounts & Import, klik 'make default' di samping email domain agar otomatis terpilih saat kirim email baru."),
        ("4. Balas Otomatis (Reply Match)", "Pilih opsi 'Reply from the same address the message was sent to' agar saat membalas email klien selalu otomatis menggunakan email domain.")
    ]
    for title, desc in usage_tips:
        p = tf_u1.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEAL_DARK
        p.space_after = Pt(1)
        
        p_sub = tf_u1.add_paragraph()
        p_sub.text = desc
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = C_TEXT_MUTED
        p_sub.space_after = Pt(6)

    # Right Card: Penggunaan di Mobile App (Android / iOS)
    add_card(s7, 6.8, 1.6, 5.7, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_u2 = s7.shapes.add_textbox(Inches(7.1), Inches(1.8), Inches(5.1), Inches(4.8))
    tf_u2 = tb_u2.text_frame
    tf_u2.word_wrap = True
    
    p = tf_u2.paragraphs[0]
    p.text = "Akses Melalui Smartphone (Android & iOS)"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(12)
    
    mobile_tips = [
        ("Otomatis Tersinkronisasi", "Setelah disetel di Web Gmail, email domain otomatis bisa dipakai mengirim & menerima di aplikasi Gmail HP tanpa instalasi tambahan."),
        ("Notifikasi Real-time", "Notifikasi pesan masuk langsung muncul di HP via push notification Gmail."),
        ("Alternatif Webmail Langsung", "Bisa juga dibuka langsung via browser di https://mail.nattuglobalsynergy.co.id:2096 menggunakan username & password."),
        ("Keamanan Multi-Factor (2FA)", "Disarankan mengaktifkan 2-Step Verification pada akun Gmail utama untuk perlindungan maksimal.")
    ]
    for title, desc in mobile_tips:
        p = tf_u2.add_paragraph()
        p.text = f"📱 {title}"
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = RGBColor(79, 70, 229)
        p.space_after = Pt(1)
        
        p_sub = tf_u2.add_paragraph()
        p_sub.text = desc
        p_sub.font.size = Pt(10.5)
        p_sub.font.color.rgb = C_TEXT_MUTED
        p_sub.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 8: SOP PENGAJUAN TERKAIT EMAIL
    # -------------------------------------------------------------
    s8 = prs.slides.add_slide(blank_slide_layout)
    add_header(s8, "SOP Pengajuan Email: Pembuatan Baru, Reset & Kendala", "TATA KELOLA LAYANAN")
    
    # 3 Service Cards
    sop_email = [
        ("Penambahan Akun Baru", [
            "Tujuan: Karyawan baru atau departemen baru.",
            "Format Pengajuan:",
            "• Nama Lengkap & Divisi",
            "• Usulan Username (misal: budi@...)",
            "• Alamat email pribadi untuk setup awal",
            "Estimasi Pengerjaan: Maksimal 1x24 Jam kerja."
        ], RGBColor(236, 253, 245), C_TEAL_DARK),
        ("Reset Password / Lupa Sandi", [
            "Tujuan: Kehilangan akses / penggantian berkala.",
            "Format Pengajuan:",
            "• Alamat email yang bermasalah",
            "• Konfirmasi PIC / Manager terkait",
            "• Password sementara akan dikirimkan via WA/Email aman",
            "Estimasi Pengerjaan: Maksimal 2-4 Jam kerja."
        ], RGBColor(238, 242, 255), RGBColor(79, 70, 229)),
        ("Penanganan Kendala / Troubleshooting", [
            "Tujuan: Email gagal kirim (bounce), error POP3/SMTP.",
            "Format Pengajuan:",
            "• Screenshot pesan error di Gmail",
            "• Alamat email pengirim & penerima",
            "• Waktu kejadian kendala",
            "Estimasi Pengerjaan: 1-6 Jam (sesuai tingkat urgensi)."
        ], RGBColor(254, 242, 242), RGBColor(220, 38, 38)),
    ]
    
    for i, (title, items, bg_c, acc_c) in enumerate(sop_email):
        x = 0.8 + (i * 3.95)
        add_card(s8, x, 1.6, 3.8, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
        
        # Header
        header_shape = s8.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x + 0.15), Inches(1.75), Inches(3.5), Inches(0.65))
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = bg_c
        header_shape.line.fill.background()
        tf = header_shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = acc_c
        p.alignment = PP_ALIGN.CENTER
        
        # Content
        tb = s8.shapes.add_textbox(Inches(x + 0.2), Inches(2.6), Inches(3.4), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = item
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT_MUTED
            p.space_after = Pt(6)

    # -------------------------------------------------------------
    # SLIDE 9: SOP PENGAJUAN UPDATE WEBSITE & BUG FIXING
    # -------------------------------------------------------------
    s9 = prs.slides.add_slide(blank_slide_layout)
    add_header(s9, "SOP Pengajuan Update Website & Bug Fixing", "TATA KELOLA LAYANAN")
    
    # Left: Kategori Permintaan & SLA
    add_card(s9, 0.8, 1.6, 6.2, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_req = s9.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_req = tb_req.text_frame
    tf_req.word_wrap = True
    
    p = tf_req.paragraphs[0]
    p.text = "Kategori Permintaan & SLA Penyelesaian"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(12)
    
    cats = [
        ("1. Minor Update (Konten & Teks)", "Perubahan nomor telepon, alamat, struktur visi misi, penambahan foto galeri divisi. SLA: 1 - 2 Hari Kerja."),
        ("2. Bug Fixing & Insiden Teknis", "Form kontak tidak terkirim, tombol tidak berfungsi, layout error di perangkat tertentu. SLA: 2 - 8 Jam Kerja."),
        ("3. Major Update / Penambahan Fitur", "Penambahan halaman baru, integrasi sistem backend baru, atau redesain modul. SLA: Estimasi diberikan via Quotation/Scope Review.")
    ]
    for title, desc in cats:
        p = tf_req.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = C_TEAL_DARK
        p.space_after = Pt(2)
        
        p_sub = tf_req.add_paragraph()
        p_sub.text = desc
        p_sub.font.size = Pt(11)
        p_sub.font.color.rgb = C_TEXT_MUTED
        p_sub.space_after = Pt(8)

    # Right: Format Pengajuan & Kanal Komunikasi
    add_card(s9, 7.3, 1.6, 5.2, 5.2, C_CARD_LIGHT, C_CARD_BORDER)
    tb_fmt = s9.shapes.add_textbox(Inches(7.5), Inches(1.8), Inches(4.8), Inches(4.8))
    tf_fmt = tb_fmt.text_frame
    tf_fmt.word_wrap = True
    
    p = tf_fmt.paragraphs[0]
    p.text = "Format Template Tiket Permintaan"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_DARK
    p.space_after = Pt(12)
    
    fmt_lines = [
        "Kirim permintaan ke WhatsApp Group Resmi / Email Support:",
        "",
        "📋 Template Format Tiket:",
        "• Jenis Request: [Update Konten / Bug Fixing]",
        "• Halaman / URL Terkait: (misal: seksi Tentang Kami)",
        "• Detail Permintaan: (Jelaskan teks lama vs teks baru)",
        "• Lampiran: (Screenshot error / File foto pengganti)",
        "• Tingkat Prioritas: [Normal / Urgent]",
        "",
        "💡 Catatan: Setiap perubahan teks disarankan sudah final (disetujui manajemen) untuk efisiensi pengerjaan."
    ]
    for line in fmt_lines:
        p = tf_fmt.add_paragraph()
        p.text = line
        if line.startswith("📋") or line.startswith("Kirim"):
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = C_TEXT_DARK
        elif line.startswith("💡"):
            p.font.italic = True
            p.font.size = Pt(10)
            p.font.color.rgb = C_AMBER
        else:
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_TEXT_MUTED
        p.space_after = Pt(3)

    # -------------------------------------------------------------
    # SLIDE 10: BEST PRACTICES KEAMANAN & MAINTENANCE
    # -------------------------------------------------------------
    s10 = prs.slides.add_slide(blank_slide_layout)
    add_header(s10, "Rekomendasi Maintenance & Keamanan Digital", "BEST PRACTICES")
    
    best_practices = [
        ("1. Keamanan Akun Email", [
            "• Gunakan kombinasi password minimal 10 karakter (huruf besar, kecil, angka, simbol).",
            "• Aktifkan 2-Factor Authentication (2FA) di akun Gmail utama.",
            "• Jangan bagikan kredensial login kepada pihak di luar perusahaan."
        ]),
        ("2. Perlindungan dari Phishing", [
            "• Waspadai email mencurigakan yang meminta transfer uang atau login mendadak.",
            "• Periksa alamat email pengirim asli secara teliti sebelum membuka lampiran file."
        ]),
        ("3. Manajemen Domain & Hosting", [
            "• Pastikan jadwal perpanjangan (renewal) tahunan domain .co.id dan hosting Biznet Gio terdaftar.",
            "• Lakukan monitoring berkala pada kapasitas penyimpanan email cPanel."
        ]),
        ("4. Standar Tanda Tangan (Signature)", [
            "• Gunakan format email signature standar PT Nattu Global Synergy yang mencakup nama, jabatan, alamat kantor, dan link website resmi."
        ])
    ]
    
    for i, (title, items) in enumerate(best_practices):
        col = i % 2
        row = i // 2
        x = 0.8 + (col * 5.95)
        y = 1.6 + (row * 2.65)
        
        add_card(s10, x, y, 5.75, 2.45, C_CARD_LIGHT, C_CARD_BORDER)
        
        tb = s10.shapes.add_textbox(Inches(x + 0.2), Inches(y + 0.2), Inches(5.35), Inches(2.1))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = C_TEAL_DARK
        p.space_after = Pt(6)
        
        for item in items:
            p = tf.add_paragraph()
            p.text = item
            p.font.size = Pt(10.5)
            p.font.color.rgb = C_TEXT_MUTED
            p.space_after = Pt(4)

    # -------------------------------------------------------------
    # SLIDE 11: PENUTUP & SESI TANYA JAWAB (Dark Theme)
    # -------------------------------------------------------------
    s11 = prs.slides.add_slide(blank_slide_layout)
    bg11 = s11.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg11.fill.solid()
    bg11.fill.fore_color.rgb = C_BG_DARK
    bg11.line.fill.background()
    
    dec11 = s11.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.0), Inches(11.73), Inches(5.5))
    dec11.fill.solid()
    dec11.fill.fore_color.rgb = C_BG_DARK_CARD
    dec11.line.color.rgb = C_TEAL_DARK
    dec11.line.width = Pt(1.5)
    
    tb_close = s11.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(10.5), Inches(4.0))
    tf_cl = tb_close.text_frame
    tf_cl.word_wrap = True
    
    p = tf_cl.paragraphs[0]
    p.text = "TERIMA KASIH"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_LIGHT
    p.space_after = Pt(10)
    
    p = tf_cl.add_paragraph()
    p.text = "Sesi Tanya Jawab (Q&A) & Praktik Bersama Setup Email"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = C_TEAL_ACCENT
    p.space_after = Pt(20)
    
    p = tf_cl.add_paragraph()
    p.text = "Kontak Dukungan Teknis & Layanan:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = C_TEXT_LIGHT
    p.space_after = Pt(6)
    
    support_info = [
        "🌐 Website: https://www.nattuglobalsynergy.co.id",
        "✉️ Email PIC: info@nattuglobalsynergy.co.id",
        "💬 WhatsApp Support Group: NSG Project & Maintenance",
        "🕒 Jam Layanan Support: Senin - Jumat (09.00 - 17.00 WIB)"
    ]
    for s in support_info:
        p = tf_cl.add_paragraph()
        p.text = s
        p.font.size = Pt(12)
        p.font.color.rgb = C_TEXT_LIGHT_MUTED
        p.space_after = Pt(4)

    output_file = "/Users/dendyaditya/Projects/windy_project/Presentasi_Handover_Website_dan_Pelatihan_Email_NSG.pptx"
    prs.save(output_file)
    print(f"SUCCESS: Presentation saved to {output_file}")

if __name__ == "__main__":
    create_deck()
